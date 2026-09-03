"""razzle.render — the python-pptx render core: a deck spec + a master + figures/logos → .pptx.

The master owns the look. razzle clones the master's layouts by ROLE (from a layout descriptor),
fills the text placeholders, and places figures/logos as pictures fitted into the descriptor-named
placeholder boxes — python-pptx cannot insert a picture into an OBJECT placeholder, so we `add_picture`
at the placeholder's geometry and drop the empty placeholder. Deterministic; the deck spec is the
durable artifact, the `.pptx` the output the author polishes.

A deck spec is a list of slides:
    {"role": "title", "title": "...", "subtitle": "..."}
    {"role": "figure", "title": "...", "figure": <figure-id>, "citation": "..."}
    {"role": "split",   "title": "...", "body": [...], "figure": <figure-id>, "citation": "..."}
    {"role": "content", "title": "...", "body": ["bullet", "bullet"]}

A deck carries no SPEAKER notes — what does not fit on the slide is spoken, not written. The notes
pane is used for one other thing entirely: a slide with no figure may carry an `illustration`, a one
line brief for a picture that does not exist yet. That is a production TODO addressed to whoever
builds the artwork, not a script addressed to the speaker, which is why it may live there when a
sentence may not.
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# The running strips (footer/contact/venue) are ONE line in a fixed box, so they are the one place
# text can overflow silently — the box does not grow and PowerPoint does not wrap what cannot wrap.
# A long contact address is the standing example: 44 characters into a 3" strip runs off the slide.
_RUNNING_SLOTS = {"footer", "contact", "venue"}
_RUNNING_PT = 12.0        # the largest we assume the master sets these at
_MIN_RUNNING_PT = 8.0     # below this it is not a contact address, it is a smudge
_ADVANCE = 0.5            # mean glyph advance as a fraction of point size, proportional sans


def _clear_slides(prs) -> None:
    """Start from the master's layouts + theme, not its example slides. Drops each slide's
    RELATIONSHIP too (not just the sldId), so the orphaned slide parts are not re-serialised — a bare
    sldId removal leaves duplicate slide parts that corrupt the deck."""
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(qn("r:id")))
        lst.remove(sid)


def _set_text(ph, value) -> None:
    if isinstance(value, (list, tuple)):
        tf = ph.text_frame
        tf.clear()
        for i, line in enumerate(value):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(line)
    else:
        ph.text = str(value)


def _fit_running_text(ph, text: str) -> None:
    """Shrink a one-line running strip until it fits its box — and only then.

    We cannot measure text without a font stack, so we estimate (mean advance ~0.5 em) and act only
    when the estimate says it overflows: the master's own styling is left alone in the common case,
    and the overflow case gets an explicit size rather than a line running off the slide edge.
    """
    box_pt = ph.width / 12700 if ph.width else 0
    need = _ADVANCE * len(text)             # width per point of font size
    if not (box_pt and need):
        return
    fits_at = box_pt / need
    if fits_at >= _RUNNING_PT:              # the master's size already fits — do not restyle
        return
    size = Pt(max(_MIN_RUNNING_PT, round(fits_at, 1)))
    tf = ph.text_frame
    tf.word_wrap = False
    for para in tf.paragraphs:
        para.font.size = size
        for run in para.runs:
            run.font.size = size


def _strip_unused(slide, filled_idxs: set) -> None:
    """Remove the master's placeholders this slide didn't fill, so they don't render as "click to
    add text". Only the UNFILLED ones go: a master's footer and contact strips are furniture the
    descriptor is expected to fill, not litter to sweep up.

    There is deliberately no slide-number branch here. python-pptx treats DATE/FOOTER/SLIDE_NUMBER
    as *latent* placeholders and never clones them onto a new slide, so a guard preserving one
    could never fire — the number has to be ADDED (see `_add_slide_number`), not kept.
    """
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in filled_idxs:
            continue
        ph._element.getparent().remove(ph._element)


def _next_shape_id(slide) -> int:
    ids = [int(e.get("id")) for e in slide.shapes._spTree.iter(qn("p:cNvPr"))
           if (e.get("id") or "").isdigit()]
    return max(ids, default=1) + 1


def _add_slide_number(slide, layout) -> None:
    """Give the slide the master's slide-number placeholder.

    `add_slide` skips it (latent), so we deep-copy the layout's own <p:sp> onto the slide. It
    carries a <a:fld type="slidenum"> field, so PowerPoint numbers it live and the number stays
    right when slides are reordered — which a literal string would not.

    The copy arrives carrying the LAYOUT's shape id, which collides with a shape already on the
    slide (a split slide ends up with two shapes numbered 7). Shape ids must be unique within a
    slide's spTree, so the clone is renumbered before it is appended.
    """
    for ph in layout.placeholders:
        if ph.element.ph_type == "sldNum" or "SLIDE_NUMBER" in str(ph.placeholder_format.type):
            el = copy.deepcopy(ph._element)
            for cnv in el.iter(qn("p:cNvPr")):
                cnv.set("id", str(_next_shape_id(slide)))
            slide.shapes._spTree.append(el)
            return


def _set_illustration_note(slide, brief: str) -> None:
    """The ONE thing the notes pane carries: a brief for a picture this slide wants and does not
    have. Nothing a speaker would read aloud ever goes here."""
    slide.notes_slide.notes_text_frame.text = f"ILLUSTRATION: {brief}"


def _place_picture(slide, ph, img: Path) -> None:
    """Add a picture fitted (aspect-preserved) inside the placeholder box, centred, then remove the
    now-empty placeholder so it doesn't render 'click to add'."""
    left, top, bw, bh = ph.left, ph.top, ph.width, ph.height
    pic = slide.shapes.add_picture(str(img), left, top, width=bw)
    if pic.height > bh:                       # too tall for the box — refit by height
        pic._element.getparent().remove(pic._element)
        pic = slide.shapes.add_picture(str(img), left, top, height=bh)
    pic.left = left + (bw - pic.width) // 2   # centre in the box
    pic.top = top + (bh - pic.height) // 2
    ph._element.getparent().remove(ph._element)


def _logo_items(logos) -> list[dict]:
    """Accept either bare image paths or {name, logo} entries, and normalise to the latter — an
    entry with no logo is a name that will be SET IN TEXT rather than dropped."""
    out = []
    for item in logos or []:
        if isinstance(item, dict):
            out.append({"name": item.get("name", ""), "logo": item.get("logo")})
        elif item:
            out.append({"name": "", "logo": item})
    return out


def _place_logo_strip(slide, box: dict, items: list[dict]) -> None:
    """Lay every logo out in a row inside an explicit box (inches), fitted to its height and centred
    as a group; a name with no registered logo is set as text in the same row.

    A single placeholder can hold ONE picture, which is why the old `logos: [idx]` mapping could
    never show a second affiliation — it silently kept the first and dropped the rest. A strip is
    the shape the thing actually is: a row of marks whose count is not known until the interview
    has run.
    """
    if not items:
        return
    left, top = Inches(box.get("left", 0)), Inches(box.get("top", 0))
    width, height = Inches(box.get("width", 1)), Inches(box.get("height", 0.4))
    gap = Inches(box.get("gap", 0.25))

    placed = []
    for it in items:
        if it.get("logo"):
            pic = slide.shapes.add_picture(str(it["logo"]), left, top, height=height)
            placed.append(pic)
        elif it.get("name"):
            tb = slide.shapes.add_textbox(left, top, Inches(2), height)
            tf = tb.text_frame
            tf.word_wrap = False
            tf.text = it["name"]
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(10)
            tb.width = max(Inches(0.5), Inches(0.09 * len(it["name"])))   # rough text advance
            placed.append(tb)

    total = sum(sh.width for sh in placed) + gap * max(0, len(placed) - 1)
    x = left + max(0, (width - total) // 2)          # centre the row in its box
    for sh in placed:
        sh.left = x
        sh.top = top + (height - sh.height) // 2     # vertical middle, whatever each one's height
        x += sh.width + gap


def render_deck(spec: list[dict], master: str, descriptor: dict, out_path: Path, *,
                figures: dict | None = None, logos: list | None = None,
                furniture: dict | None = None) -> Path:
    """Render the deck spec onto the branded master.

    `figures` maps a slide's figure-id → an image path. `logos` is the ordered list of logo image
    paths (or {name, logo} entries) for a role's logo slots or its `logo_strip`. `furniture` is the
    deck-level running text — venue/date, the running footer, the contact address — which is the
    same on every slide and so is NOT in the spec: the composer never sees it and cannot invent it.
    A text slot takes the slide's own value first and falls back to the furniture.

    A slide with an `illustration` and no figure gets that brief in its notes pane — the only thing
    notes ever hold. The opening slide carries no page number, as a title page never does.

    A missing figure or logo is simply skipped (the box stays empty) — never a crash.
    """
    prs = Presentation(str(master))
    _clear_slides(prs)
    roles = descriptor.get("roles", {})
    figures = figures or {}
    furniture = furniture or {}
    items = _logo_items(logos)
    for slide in spec:
        rdef = roles.get(slide.get("role", "figure"))
        if rdef is None:
            continue
        layout = prs.slide_layouts[rdef["layout"]]
        s = prs.slides.add_slide(layout)
        phs = {ph.placeholder_format.idx: ph for ph in s.placeholders}
        filled: set = set()
        for slot, idx in (rdef.get("text") or {}).items():
            val = slide.get(slot) or furniture.get(slot)     # the slide's own, else the deck's
            if val and idx in phs:
                _set_text(phs[idx], val)
                if slot in _RUNNING_SLOTS and isinstance(val, str):
                    _fit_running_text(phs[idx], val)
                filled.add(idx)
        for slot, idx in (rdef.get("picture") or {}).items():
            img = figures.get(slide.get(slot))
            if img and idx in phs:
                _place_picture(s, phs[idx], Path(img))
                filled.add(idx)
        for i, idx in enumerate(rdef.get("logos") or []):      # one placeholder, one logo
            if i < len(items) and items[i].get("logo") and idx in phs:
                _place_picture(s, phs[idx], Path(items[i]["logo"]))
                filled.add(idx)
        _strip_unused(s, filled)      # drop the empty master placeholders this slide didn't use
        if rdef.get("logo_strip"):    # after the strip, so it is not swept up as unfilled
            _place_logo_strip(s, rdef["logo_strip"], items)
        if slide.get("illustration") and not slide.get("figure"):
            _set_illustration_note(s, str(slide["illustration"]))
        if slide.get("role") != "title":   # the opening slide is never numbered
            _add_slide_number(s, layout)
    prs.save(str(out_path))
    return out_path
