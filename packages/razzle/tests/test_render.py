"""razzle.render — the python-pptx render core.

A deck spec + the neutral house master + a figure → a .pptx: layouts cloned by role, text
placeholders filled, the figure placed as a picture in the content area. Logos are wired but skipped
when absent. The real master/logos live neutral (~/.config/haarpi/razzle/) and
never in the repo, so the render test is skip-guarded on their presence.
"""

from __future__ import annotations

import importlib.util
import shutil

import pytest

from razzle import assets, render

_HAS_DOT = shutil.which("dot") is not None
# a figure reaches a slide as a PNG, so a rasteriser is as much a prerequisite as graphviz — without
# one `figure.resolve(..., "png")` returns None and the picture assertion below fails for a reason
# that has nothing to do with razzle
_HAS_RASTER = importlib.util.find_spec("cairosvg") is not None or shutil.which("rsvg-convert")
_MASTER = assets.master_pptx("default")
_DESC = assets.descriptor("default")


def test_example_descriptor_documents_the_format():
    import yaml
    from importlib.resources import files
    d = yaml.safe_load((files("razzle") / "layouts" / "example.yaml").read_text())
    assert d["size"] == "16:9"
    assert d["roles"]["figure"]["picture"] == {"figure": 1}
    # idx14 on the title layout is the VENUE/DATE line, not a logo strip. The title slide carries
    # NO marks at all: the logos get their own explicit box on the acknowledgements slide, because
    # a placeholder holds exactly one picture and there may be several marks.
    assert d["roles"]["title"]["text"]["venue"] == 14
    assert "logos" not in d["roles"]["title"] and "logo_strip" not in d["roles"]["title"]
    ack = d["roles"]["acknowledgements"]
    assert set(ack["logo_strip"]) >= {"left", "top", "width", "height"}
    assert {"footer", "contact"} <= set(ack["text"])
    # the running furniture is mapped on every non-title role, so it is filled instead of stripped
    for role in ("figure", "content", "split"):
        assert {"footer", "contact"} <= set(d["roles"][role]["text"])
    # `split` puts a point beside its evidence. Its indices are NOT layout 1's: idx13 is the
    # right-hand figure box and idx14 the caption, the reverse of the figure role above.
    split = d["roles"]["split"]
    assert split["layout"] == 2
    assert split["picture"] == {"figure": 13}
    assert split["text"] == {"title": 0, "body": 1, "citation": 14, "footer": 15, "contact": 16}


def test_logos_for_is_a_list_from_the_registries():
    got = assets.logos_for(affiliations=["Nowhere University"], funders=[])
    assert isinstance(got, list)                       # unmatched → skipped, never a crash


@pytest.mark.skipif(not (_MASTER and _DESC and _HAS_DOT and _HAS_RASTER),
                    reason="neutral house master, graphviz or SVG→PNG rasteriser absent")
def test_render_deck_against_the_house_master(tmp_path):
    from haarpi import figure, project
    figure.write_figure(tmp_path, "demo", figure.stage_dag(project.DEFAULT_STAGES))
    png = figure.resolve(tmp_path, "demo", "stageLadder", "png", width=1200)

    spec = [
        {"role": "title", "title": "A research talk", "subtitle": "the through-line"},
        {"role": "figure", "title": "The pipeline", "figure": "ladder", "citation": "[Ref 1]"},
        {"role": "split", "title": "The point, beside its evidence", "figure": "ladder",
         "body": ["it composes", "it's grounded"], "citation": "[Ref 2]"},
        {"role": "content", "title": "Takeaways", "body": ["it composes", "it's grounded"]},
    ]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "deck.pptx",
                             figures={"ladder": png})

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 4
    assert any("A research talk" in (sh.text_frame.text if sh.has_text_frame else "")
               for sh in prs.slides[0].shapes)
    assert any(sh.shape_type == 13 for sh in prs.slides[1].shapes)     # 13 == PICTURE
    # the split slide carries BOTH its bullets and its figure — the whole reason the role exists
    texts = " ".join(sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame)
    assert "it composes" in texts and "[Ref 2]" in texts
    assert any(sh.shape_type == 13 for sh in prs.slides[2].shapes)


def test_a_deck_carries_no_speaker_notes(tmp_path):
    """Notes are where the essay goes when the slide refuses it. A `notes` key on a slide must not
    reach the .pptx even if something upstream puts one there."""
    spec = [{"role": "content", "title": "T", "body": ["a"], "notes": "an essay"}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx")

    from pptx import Presentation
    prs = Presentation(str(out))
    assert not prs.slides[0].has_notes_slide


# ------------------------------------------------------- the master's running furniture
def _texts(slide):
    return " | ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_deck_level_furniture_fills_the_strips_instead_of_being_stripped(tmp_path):
    """venue/footer/contact are the same on every slide, so they are NOT in the spec — they come
    from `furniture`. Before this they were mapped nowhere and _strip_unused deleted the master's
    footer and contact strips off every slide."""
    spec = [{"role": "title", "title": "A talk", "subtitle": "Ada, Bo"},
            {"role": "content", "title": "Points", "body": ["a"]}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx",
                             furniture={"venue": "CSS2026 · 31 Oct", "footer": "CSS2026 | A talk",
                                        "contact": "ada@ucb.edu"})
    from pptx import Presentation
    prs = Presentation(str(out))
    assert "CSS2026 · 31 Oct" in _texts(prs.slides[0])          # the title slide's venue/date line
    assert "CSS2026 | A talk" in _texts(prs.slides[1])          # the running footer
    assert "ada@ucb.edu" in _texts(prs.slides[1])               # the running contact


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_a_slide_value_beats_the_deck_furniture(tmp_path):
    spec = [{"role": "content", "title": "T", "body": ["a"], "footer": "this slide only"}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx",
                             furniture={"footer": "the deck's"})
    from pptx import Presentation
    assert "this slide only" in _texts(Presentation(str(out)).slides[0])


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_every_slide_gets_the_masters_slide_number(tmp_path):
    """python-pptx never CLONES a slide-number placeholder (it is latent), so it has to be added.
    The clone carries the master's <a:fld type="slidenum">, which numbers live rather than baking
    a string that would go wrong the moment slides are reordered."""
    spec = [{"role": "content", "title": str(i), "body": ["x"]} for i in range(3)]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx")
    from pptx import Presentation
    from pptx.oxml.ns import qn
    for slide in Presentation(str(out)).slides:
        flds = slide.shapes._spTree.findall(f".//{qn('a:fld')}")
        assert any(f.get("type") == "slidenum" for f in flds)


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_the_opening_slide_is_not_numbered(tmp_path):
    """A title page carries no page number. Every slide after it does."""
    spec = [{"role": "title", "title": "A talk", "subtitle": "Ada"},
            {"role": "content", "title": "Points", "body": ["a"]}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx")
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs = Presentation(str(out))

    def _numbered(slide):
        return any(f.get("type") == "slidenum"
                   for f in slide.shapes._spTree.findall(f".//{qn('a:fld')}"))

    assert not _numbered(prs.slides[0])
    assert _numbered(prs.slides[1])


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_the_slide_number_clone_gets_a_fresh_shape_id(tmp_path):
    """The clone arrives carrying the LAYOUT's shape id, which collides with a shape already on the
    slide — a split slide ended up with two shapes numbered 7. Ids must be unique per spTree."""
    spec = [{"role": "split", "title": "T", "body": ["a"], "citation": "[R]"},
            {"role": "content", "title": "U", "body": ["b"]}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx")
    from pptx import Presentation
    from pptx.oxml.ns import qn
    for slide in Presentation(str(out)).slides:
        ids = [e.get("id") for e in slide.shapes._spTree.iter(qn("p:cNvPr"))]
        assert len(ids) == len(set(ids)), ids


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_a_long_contact_address_is_shrunk_to_fit_its_strip(tmp_path):
    """The running strips are one line in a fixed box, so a 44-character address ran off the slide.
    It is shrunk only when it overflows — a short one keeps the master's own styling."""
    long_ = "d.cale.reeves@computationalsocialscience.org"
    out = render.render_deck([{"role": "content", "title": "T", "body": ["a"]}],
                             _DESC["master_path"], _DESC, tmp_path / "long.pptx",
                             furniture={"contact": long_, "footer": "CSS2026 | A talk"})
    short = render.render_deck([{"role": "content", "title": "T", "body": ["a"]}],
                               _DESC["master_path"], _DESC, tmp_path / "short.pptx",
                               furniture={"contact": "ada@ucb.edu", "footer": "CSS2026 | A talk"})
    from pptx import Presentation

    def _size_of(path, needle):
        for sh in Presentation(str(path)).slides[0].shapes:
            if sh.has_text_frame and needle in sh.text_frame.text:
                return sh.text_frame.paragraphs[0].font.size
        raise AssertionError(f"{needle} not on the slide")

    sized = _size_of(out, long_)
    assert sized is not None and sized.pt <= 10        # shrunk to fit the ~3" strip
    assert _size_of(short, "ada@ucb.edu") is None      # untouched — the master's size already fits


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_an_illustration_brief_is_the_only_thing_in_the_notes(tmp_path):
    """Notes carry no speech — but a slide with nothing to show may brief the picture it wants,
    addressed to whoever draws it. A slide that already has a figure gets no brief."""
    spec = [{"role": "content", "title": "T", "body": ["a"],
             "illustration": "a piano roll with the target phrase above a scrambled one"},
            {"role": "content", "title": "U", "body": ["b"]}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert prs.slides[0].notes_slide.notes_text_frame.text == (
        "ILLUSTRATION: a piano roll with the target phrase above a scrambled one")
    assert not prs.slides[1].has_notes_slide


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_the_logo_strip_shows_every_mark_not_just_the_first(tmp_path):
    """A placeholder holds ONE picture, which is why `logos: [idx]` could only ever show the first
    affiliation. The strip lays out as many as the interview chose."""
    logos = sorted((assets.home() / "logos").glob("*.png"))[:2]
    if len(logos) < 2:
        pytest.skip("need two logo files in the neutral home")
    spec = [{"role": "acknowledgements", "title": "Acknowledgements"}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx", logos=logos)
    from pptx import Presentation
    pics = [sh for sh in Presentation(str(out)).slides[0].shapes if sh.shape_type == 13]
    assert len(pics) == 2
    assert pics[0].left < pics[1].left          # laid out in a row, in order


@pytest.mark.skipif(not (_MASTER and _DESC), reason="neutral house master absent")
def test_an_unregistered_affiliation_is_set_in_text_not_dropped(tmp_path):
    """The registries have always promised "degrades to text (name only)" and never done it — an
    affiliation the author picked in the interview just vanished."""
    spec = [{"role": "acknowledgements", "title": "Acknowledgements"}]
    out = render.render_deck(spec, _DESC["master_path"], _DESC, tmp_path / "d.pptx",
                             logos=[{"name": "Nowhere University", "logo": None}])
    from pptx import Presentation
    assert "Nowhere University" in _texts(Presentation(str(out)).slides[0])


# ------------------------------------------------------- registry name matching
def test_an_alias_reaches_the_logo_filed_under_a_short_name(tmp_path, monkeypatch):
    """Manifests write the official long name; the registry files the logo under the short one."""
    home = tmp_path / "home"
    (home / "logos").mkdir(parents=True)
    (home / "logos" / "v.png").write_bytes(b"x")
    (home / "affiliations.yaml").write_text(
        'VCUarts Qatar:\n  logo: logos/v.png\n  aliases:\n    - "VCU School of the Arts, Qatar"\n')
    monkeypatch.setenv("RAZZLE_HOME", str(home))

    assert assets.logos_for(["VCUarts Qatar"])[0].name == "v.png"          # the key
    assert assets.logos_for(["VCU School of the Arts, Qatar"])[0].name == "v.png"   # an alias
    assert assets.logos_for(["  vcuarts   qatar  "])[0].name == "v.png"    # case/whitespace
    # an unmatched name keeps its place, with no logo, so it can be set in text
    entries = assets.logo_entries(["Nowhere University"])
    assert entries == [{"name": "Nowhere University", "logo": None}]
