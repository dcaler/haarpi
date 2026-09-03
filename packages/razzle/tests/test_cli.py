"""razzle CLI — `razzle deck` (gather + author) and `razzle render` (spec → .pptx)."""

from __future__ import annotations

import json

import pytest

from haarpi import figure, project
from razzle import assets, cli

_HAS_MASTER = assets.master_pptx("default") is not None


def _fixture(root):
    project.save_manifest(project.Manifest(name="demo", short_title="demo", brief="x"), root)
    figure.write_figure(root, "demo", figure.stage_dag(project.DEFAULT_STAGES), render_svg=False)


def test_deck_no_launch_gathers_and_scaffolds(tmp_path, capsys):
    _fixture(tmp_path)
    rc = cli.main(["deck", "--dir", str(tmp_path), "--format", "shorttalk", "--no-launch"])
    assert rc == 0
    out = capsys.readouterr().out
    assert "shorttalk" in out and "~11 slides" in out          # sized to the format
    assert (tmp_path / "slides" / "shorttalk").is_dir()


def test_deck_rejects_unknown_format(tmp_path):
    _fixture(tmp_path)
    assert cli.main(["deck", "--dir", str(tmp_path), "--format", "keynote", "--no-launch"]) == 2


@pytest.mark.skipif(not _HAS_MASTER, reason="neutral house master absent")
def test_render_builds_pptx_from_a_spec(tmp_path):
    _fixture(tmp_path)
    fmt_dir = tmp_path / "slides" / "shorttalk"
    fmt_dir.mkdir(parents=True)
    (fmt_dir / "spec.json").write_text(json.dumps([
        {"role": "title", "title": "T", "subtitle": "S"},
        {"role": "figure", "title": "F", "figure": "stageLadder"}]))
    assert cli.main(["render", "--dir", str(tmp_path), "--format", "shorttalk"]) == 0
    from pptx import Presentation
    from haarpi import naming
    pptx = fmt_dir / naming.major_name("demo", "pptx", infix="deck")   # {date}_demo_deck_ra.pptx
    prs = Presentation(str(pptx))
    # 2 authored + the acknowledgements slide razzle appends itself
    assert pptx.is_file() and len(prs.slides) == 3
    assert any("Acknowledgements" in (sh.text_frame.text if sh.has_text_frame else "")
               for sh in prs.slides[-1].shapes)


def test_render_without_a_spec_errors(tmp_path):
    _fixture(tmp_path)
    assert cli.main(["render", "--dir", str(tmp_path), "--format", "longtalk"]) == 1


def test_claude_headless_fails_loudly_when_there_is_no_claude(tmp_path, monkeypatch, capsys):
    """A queued task that exits 0 without authoring is worse than one that fails: the board would
    chain the review of a deck that was never written. (--headless only qualifies --claude; the
    default path is the local brain and needs no binary on PATH.)"""
    from razzle import cli
    monkeypatch.setattr(cli.shutil, "which", lambda _n: None)
    rc = cli.main(["deck", "--format", "longtalk", "--claude", "--headless", "--dir", str(tmp_path)])
    assert rc == 1
    assert "needs `claude` on PATH" in capsys.readouterr().err


def test_the_default_path_is_the_local_brain_not_a_session(tmp_path, monkeypatch, capsys):
    """`razzle deck` with no flags must never shell out to a cloud session — offline-first is the
    default, and the cloud coordinator is an explicit, human-invoked opt-in."""
    from razzle import cli
    called = {}

    def _fake_build(root, fmt, brain, *, master="default", out=None, bundle=None):
        called["fmt"], called["brain"] = fmt, brain
        return {"spec": [{"role": "title", "title": "T"}], "pptx": tmp_path / "d.pptx"}

    monkeypatch.setattr(cli, "_local_brain", lambda: "LOCAL")
    monkeypatch.setattr("razzle.deck.build_deck", _fake_build)
    monkeypatch.setattr(cli.subprocess, "run",
                        lambda *a, **k: pytest.fail("the default path shelled out to a session"))
    assert cli.main(["deck", "--format", "shorttalk", "--dir", str(tmp_path)]) == 0
    assert called == {"fmt": "shorttalk", "brain": "LOCAL"}


def test_render_enforces_the_budgets_on_a_hand_authored_spec(tmp_path, monkeypatch):
    """The spec is written by a session, so it arrives with a session's habits. Rendering it raw
    meant compose.py's rules only ever applied to a path nothing calls."""
    import json as _json
    from razzle import assets, cli
    if not assets.descriptor("default"):
        import pytest as _pytest
        _pytest.skip("neutral house master absent")

    root = tmp_path
    (root / "slides" / "longtalk").mkdir(parents=True)
    (root / "slides" / "longtalk" / "spec.json").write_text(_json.dumps([
        {"role": "title", "title": "T", "subtitle": "s"},
        {"role": "content", "title": "Too much", "body": ["a", "b", "c", "d", "e"],
         "notes": "an essay nobody reads"},
    ]))
    assert cli.main(["render", "--format", "longtalk", "--dir", str(root)]) == 0

    from pptx import Presentation
    out = next((root / "slides" / "longtalk").glob("*_deck_ra.pptx"))
    prs = Presentation(str(out))
    body = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "d" not in body.split() and "e" not in body.split()   # capped at three
    assert not prs.slides[1].has_notes_slide                      # notes dropped
